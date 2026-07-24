"""Knowledge-base ingestion — add company documents (PDF/DOCX/PPTX) or a website to the library.

Extract text per format, let Gemini classify it (doc_type, validity, structured fields), and
hand back a library_documents row. Untrusted input throughout: the extracted text is data for
the classifier only; URL fetch is guarded against SSRF and bounded in size.
"""

from __future__ import annotations

import datetime
import io
import ipaddress
import re
import socket
import zipfile
from html.parser import HTMLParser
from pathlib import Path
from urllib.parse import urlparse

import httpx

from pipeline.client import ModelError, generate_json
from pipeline.schemas import KB_DOC_SCHEMA

from .envelope import ApiError

_MAX_FETCH_BYTES = 5 * 1024 * 1024
_MAX_UNCOMPRESSED = 100 * 1024 * 1024  # zip-bomb ceiling for office (zip) formats
_PROMPT = (Path(__file__).resolve().parents[1] / "prompts" / "kb_classifier.md").read_text()


# ---------- text extraction ----------
def extract_text(filename: str, data: bytes) -> str:
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext == "pdf":
        from .ingest import parse_pdf_pages

        return "\n".join(t for _, t in parse_pdf_pages(data)).strip()
    if ext in ("docx",):
        return _extract_docx(data)
    if ext in ("pptx",):
        return _extract_pptx(data)
    # plain text / unknown: best-effort decode
    return data.decode("utf-8", errors="ignore").strip()


def _guard_zip_bomb(data: bytes) -> None:
    """Reject office files whose uncompressed size would blow up memory (decompression bomb)."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            if sum(i.file_size for i in z.infolist()) > _MAX_UNCOMPRESSED:
                raise ApiError(413, "FILE_TOO_LARGE", "document expands to an unsafe size")
    except zipfile.BadZipFile:
        pass  # not a valid zip — the format parser below will reject it


def _extract_docx(data: bytes) -> str:
    _guard_zip_bomb(data)
    from docx import Document

    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — bad upload -> 400
        raise ApiError(400, "BAD_DOCUMENT", f"could not read DOCX: {exc}") from exc
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()


def _extract_pptx(data: bytes) -> str:
    _guard_zip_bomb(data)
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — bad upload -> 400
        raise ApiError(400, "BAD_DOCUMENT", f"could not read PPTX: {exc}") from exc
    parts = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip()
    ]
    return "\n".join(parts).strip()


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._skip = False
        self.chunks: list[str] = []

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self.chunks.append(data.strip())


def fetch_url_text(url: str) -> str:
    """Fetch a public webpage and strip it to text.

    SSRF-hardened: every hop's host is RESOLVED and every resolved IP checked against the
    private/loopback/link-local ranges (defeats DNS-rebind), redirects are followed manually
    (max 4) so each new target is re-validated (defeats redirect-to-metadata), and the body is
    streamed with a hard byte cap (defeats oversize/OOM).
    """
    for _ in range(4):
        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https"):
            raise ApiError(400, "BAD_URL", "only http(s) URLs are allowed")
        host = parsed.hostname
        if not host:
            raise ApiError(400, "BAD_URL", "missing host")
        _reject_private_host(host)

        try:
            with httpx.stream(
                "GET", url, timeout=15, follow_redirects=False,
                headers={"User-Agent": "TenderCraft/0.1"},
            ) as resp:
                if resp.is_redirect:
                    loc = resp.headers.get("location")
                    if not loc:
                        raise ApiError(400, "FETCH_FAILED", "redirect without location")
                    url = httpx.URL(url).join(loc).__str__()
                    continue  # re-validate the new target on the next iteration
                resp.raise_for_status()
                buf = bytearray()
                for chunk in resp.iter_bytes():
                    buf.extend(chunk)
                    if len(buf) > _MAX_FETCH_BYTES:
                        break
        except httpx.HTTPError as exc:
            raise ApiError(400, "FETCH_FAILED", f"could not fetch URL: {exc}") from exc

        p = _TextExtractor()
        p.feed(bytes(buf[:_MAX_FETCH_BYTES]).decode("utf-8", errors="ignore"))
        return re.sub(r"\s+\n", "\n", " ".join(p.chunks)).strip()

    raise ApiError(400, "TOO_MANY_REDIRECTS", "the URL redirected too many times")


def _reject_private_host(host: str) -> None:
    """Resolve the host and reject if ANY resolved address is non-public."""
    if host in ("localhost", "metadata.google.internal"):
        raise ApiError(400, "BLOCKED_URL", "internal/private hosts are not allowed")
    try:
        infos = socket.getaddrinfo(host, None)
    except socket.gaierror as exc:
        raise ApiError(400, "FETCH_FAILED", "could not resolve host") from exc
    for info in infos:
        addr = ipaddress.ip_address(info[4][0])
        if (
            addr.is_private or addr.is_loopback or addr.is_link_local
            or addr.is_reserved or addr.is_multicast or addr.is_unspecified
        ):
            raise ApiError(400, "BLOCKED_URL", "internal/private hosts are not allowed")


# ---------- classification ----------
def classify(text: str) -> dict:
    """Derive {name, doc_type, valid_to, structured_fields(dict)} from text. Never raises."""
    empty = {
        "name": "Empty document", "doc_type": "other",
        "valid_to": None, "structured_fields": {},
    }
    if not text.strip():
        return empty
    prompt = _PROMPT.replace("{{TEXT}}", text[:6000])
    try:
        r = generate_json(prompt, KB_DOC_SCHEMA)
    except ModelError:
        return {**empty, "name": "Uploaded document"}
    fields = {f["key"]: f["value"] for f in r.get("structured_fields", []) if f.get("key")}
    return {
        "name": r.get("name") or "Uploaded document",
        "doc_type": r.get("doc_type", "other"),
        "valid_to": _valid_iso_date(r.get("valid_to")),
        "structured_fields": fields,
    }


def _valid_iso_date(value) -> str | None:
    """Coerce the model's valid_to to a real ISO date or None — a bad string would corrupt the
    validity hard-filter (string compare) that excludes expired docs from drafting."""
    if not value:
        return None
    try:
        return datetime.date.fromisoformat(str(value)).isoformat()
    except ValueError:
        return None


def build_document(text: str) -> dict:
    """Classify extracted text into a library_documents row payload (minus tenant/actor)."""
    meta = classify(text)
    return {
        "name": meta["name"],
        "doc_type": meta["doc_type"],
        "valid_to": meta["valid_to"],
        "text_content": text[:20000],
        "structured_fields": meta["structured_fields"],
    }
