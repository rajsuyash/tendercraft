"""Knowledge-base ingestion — format extraction, URL guard, classifier."""

import io

import pytest

from app import knowledge
from app.envelope import ApiError


def test_extract_plain_text():
    assert knowledge.extract_text("notes.txt", b"hello world") == "hello world"


def test_extract_docx_roundtrip():
    from docx import Document

    doc = Document()
    doc.add_paragraph("Turnover certificate for FY25")
    doc.add_paragraph("Value: 9.7 Cr")
    buf = io.BytesIO()
    doc.save(buf)
    text = knowledge.extract_text("cert.docx", buf.getvalue())
    assert "Turnover certificate" in text
    assert "9.7 Cr" in text


def test_extract_pptx_roundtrip():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Company capability deck"
    buf = io.BytesIO()
    prs.save(buf)
    text = knowledge.extract_text("deck.pptx", buf.getvalue())
    assert "Company capability deck" in text


def test_bad_docx_is_400():
    with pytest.raises(ApiError) as e:
        knowledge.extract_text("x.docx", b"not a docx")
    assert e.value.code == "BAD_DOCUMENT"


def test_url_rejects_non_http():
    with pytest.raises(ApiError) as e:
        knowledge.fetch_url_text("ftp://example.com/x")
    assert e.value.code == "BAD_URL"


@pytest.mark.parametrize("url", [
    "http://localhost/x",
    "http://127.0.0.1/x",
    "http://10.0.0.1/x",
    "http://169.254.169.254/latest/meta-data",
])
def test_url_blocks_private_and_loopback(url):
    with pytest.raises(ApiError) as e:
        knowledge.fetch_url_text(url)
    assert e.value.code == "BLOCKED_URL"


def test_url_blocks_dns_resolving_to_private(monkeypatch):
    # A public-looking name that resolves to a link-local IP (DNS-rebind / metadata bypass).
    monkeypatch.setattr(
        knowledge.socket, "getaddrinfo",
        lambda host, port: [(2, 1, 6, "", ("169.254.169.254", 0))],
    )
    with pytest.raises(ApiError) as e:
        knowledge.fetch_url_text("http://evil.example.com/x")
    assert e.value.code == "BLOCKED_URL"


class _FakeResp:
    def __init__(self, *, redirect_to=None, body=b""):
        self._redirect_to = redirect_to
        self._body = body
        self.is_redirect = redirect_to is not None
        self.headers = {"location": redirect_to} if redirect_to else {}

    def __enter__(self):
        return self

    def __exit__(self, *a):
        return False

    def raise_for_status(self):
        pass

    def iter_bytes(self):
        yield self._body


def test_url_blocks_redirect_to_metadata(monkeypatch):
    # First hop is a real public host; it 302s to the cloud metadata IP — must be re-validated.
    monkeypatch.setattr(knowledge.socket, "getaddrinfo",
                        lambda host, port: [(2, 1, 6, "", (host, 0))]
                        if host == "169.254.169.254"
                        else [(2, 1, 6, "", ("93.184.216.34", 0))])
    monkeypatch.setattr(knowledge.httpx, "stream",
                        lambda *a, **k: _FakeResp(redirect_to="http://169.254.169.254/latest"))
    with pytest.raises(ApiError) as e:
        knowledge.fetch_url_text("http://public.example.com/redir")
    assert e.value.code == "BLOCKED_URL"


def test_url_happy_path_streams_and_strips(monkeypatch):
    monkeypatch.setattr(knowledge.socket, "getaddrinfo",
                        lambda host, port: [(2, 1, 6, "", ("93.184.216.34", 0))])
    monkeypatch.setattr(knowledge.httpx, "stream",
                        lambda *a, **k: _FakeResp(body=b"<body><p>We build things.</p></body>"))
    assert "We build things." in knowledge.fetch_url_text("http://public.example.com/")


@pytest.mark.parametrize("bad", ["not-a-date", "2027-13-99", "", "soon"])
def test_valid_iso_date_rejects_garbage(bad):
    assert knowledge._valid_iso_date(bad) is None


def test_valid_iso_date_passes_real_date():
    assert knowledge._valid_iso_date("2027-09-01") == "2027-09-01"


def test_html_stripped_to_text():
    p = knowledge._TextExtractor()
    p.feed("<html><style>.a{}</style><body><h1>About Us</h1><script>x()</script><p>We build things.</p></body></html>")
    joined = " ".join(p.chunks)
    assert "About Us" in joined and "We build things." in joined
    assert "x()" not in joined  # script content dropped


def test_classify_empty_text_is_other():
    meta = knowledge.classify("   ")
    assert meta["doc_type"] == "other"
    assert meta["structured_fields"] == {}


def test_classify_model_error_falls_back(monkeypatch):
    def boom(*a, **k):
        raise knowledge.ModelError("down")

    monkeypatch.setattr(knowledge, "generate_json", boom)
    meta = knowledge.classify("some real text here")
    assert meta["doc_type"] == "other"
    assert meta["name"] == "Uploaded document"


def test_classify_maps_structured_fields(monkeypatch):
    monkeypatch.setattr(knowledge, "generate_json", lambda p, s, **k: {
        "name": "ISO 9001 Certificate",
        "doc_type": "certification",
        "valid_to": "2027-09-01",
        "structured_fields": [{"key": "cert_no", "value": "IN-9001-44821"}],
    })
    meta = knowledge.classify("ISO 9001 cert text")
    assert meta["name"] == "ISO 9001 Certificate"
    assert meta["doc_type"] == "certification"
    assert meta["valid_to"] == "2027-09-01"
    assert meta["structured_fields"] == {"cert_no": "IN-9001-44821"}


def _stub_ingest(monkeypatch):
    from app import db, knowledge_routes
    monkeypatch.setattr(knowledge_routes.knowledge, "build_document",
                        lambda t: {"name": "X", "doc_type": "completion", "valid_to": None,
                                   "text_content": t, "structured_fields": {}})
    monkeypatch.setattr(db, "insert_library_document",
                        lambda tid, doc, actor: {"id": "doc-1", "name": "X",
                                                 "doc_type": "completion", "valid_to": None})
    calls: list = []
    monkeypatch.setattr(db, "upsert_readiness_decision", lambda *a, **k: calls.append((a, k)))
    return knowledge_routes, calls


def test_ingest_links_document_to_criterion(monkeypatch):
    knowledge_routes, calls = _stub_ingest(monkeypatch)
    res = knowledge_routes._ingest_text("t1", "u1", "text", "crit-1", "tender-1")
    assert res["id"] == "doc-1"
    assert len(calls) == 1
    assert calls[0][1]["document_id"] == "doc-1"
    assert calls[0][0] == ("t1", "tender-1", "crit-1")


def test_ingest_without_criterion_does_not_link(monkeypatch):
    knowledge_routes, calls = _stub_ingest(monkeypatch)
    knowledge_routes._ingest_text("t1", "u1", "text")  # no criterion/tender
    assert calls == []


def test_build_document_shape(monkeypatch):
    monkeypatch.setattr(knowledge, "generate_json", lambda p, s, **k: {
        "name": "Completion Cert", "doc_type": "completion", "valid_to": None, "structured_fields": [],
    })
    doc = knowledge.build_document("completion certificate text")
    assert doc["name"] == "Completion Cert"
    assert doc["doc_type"] == "completion"
    assert doc["valid_to"] is None
    assert doc["text_content"].startswith("completion")
